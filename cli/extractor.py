import csv
import logging
import zipfile
from pathlib import Path

log = logging.getLogger(__name__)


def extract_text_from_pdf(filepath: Path) -> str:
    """Extract text from a PDF file."""
    import pdfplumber

    text_parts = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"--- Page {i + 1} ---\n{page_text}")
    except Exception as e:
        log.warning("Failed to extract text from %s: %s", filepath.name, e)

    return "\n\n".join(text_parts)


def extract_text_from_pptx(filepath: Path) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation

    text_parts = []
    try:
        prs = Presentation(str(filepath))
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                text_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))
    except Exception as e:
        log.warning("Failed to extract text from %s: %s", filepath.name, e)

    return "\n\n".join(text_parts)


def extract_text_from_excel(filepath: Path) -> str:
    """Extract text from an Excel file (.xlsx/.xls)."""
    import openpyxl

    text_parts = []
    try:
        wb = openpyxl.load_workbook(str(filepath), data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append("\t".join(cells))
            if rows:
                text_parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
    except Exception as e:
        log.warning("Failed to extract text from %s: %s", filepath.name, e)

    return "\n\n".join(text_parts)


def extract_text_from_csv(filepath: Path) -> str:
    """Extract text from a CSV file."""
    text_parts = []
    try:
        with open(filepath, newline="", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                if any(cell.strip() for cell in row):
                    text_parts.append("\t".join(row))
    except Exception as e:
        log.warning("Failed to extract text from %s: %s", filepath.name, e)

    return "\n".join(text_parts)


def extract_text_from_zip(filepath: Path) -> str:
    """Extract text from supported files inside a ZIP archive."""
    text_parts = []
    try:
        with zipfile.ZipFile(filepath, "r") as zf:
            for name in zf.namelist():
                if name.endswith("/"):
                    continue
                ext = Path(name).suffix.lower()
                if ext not in _ZIP_SUPPORTED_EXTENSIONS:
                    continue
                log.info("  Extracting from zip member: %s", name)
                data = zf.read(name)
                # Write to a temp file so extraction functions can read it
                import tempfile
                with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                    tmp.write(data)
                    tmp_path = Path(tmp.name)
                try:
                    member_text = extract_text(tmp_path)
                    if member_text.strip():
                        text_parts.append(f"--- {name} ---\n{member_text}")
                finally:
                    tmp_path.unlink(missing_ok=True)
    except zipfile.BadZipFile:
        log.warning("Bad zip file: %s", filepath.name)
    except Exception as e:
        log.warning("Failed to extract text from %s: %s", filepath.name, e)

    return "\n\n".join(text_parts)


_ZIP_SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".xlsx", ".xls", ".csv"}


def extract_text(filepath: Path) -> str:
    """Extract text from a supported file type."""
    ext = filepath.suffix.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(filepath)
    elif ext in (".xlsx", ".xls"):
        return extract_text_from_excel(filepath)
    elif ext == ".csv":
        return extract_text_from_csv(filepath)
    elif ext == ".zip":
        return extract_text_from_zip(filepath)
    else:
        log.warning("Unsupported file type: %s", ext)
        return ""


def extract_all(filepaths: list[Path]) -> dict[str, str]:
    """Extract text from all files. Returns {filename: text} dict."""
    results = {}
    for fp in filepaths:
        text = extract_text(fp)
        if text.strip():
            results[fp.name] = text
            log.info("  Extracted %d chars from %s", len(text), fp.name)
        else:
            log.warning("  No text extracted from %s", fp.name)
    return results
